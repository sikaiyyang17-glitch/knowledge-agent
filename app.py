from flask import Flask, render_template, request, redirect, url_for, jsonify, send_file
from flask_sqlalchemy import SQLAlchemy
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from dotenv import load_dotenv
from google import genai
from groq import Groq
import bcrypt
import os
import zipfile
import io
import json
import re
import numpy as np
from datetime import datetime

load_dotenv()

from sentence_transformers import SentenceTransformer
print("Loading embedding model...")
embedder = SentenceTransformer('all-MiniLM-L6-v2')
print("Embedding model loaded.")

import faiss
from rank_bm25 import BM25Okapi

def semantic_chunk(text, filename):
    if not text or not text.strip():
        return []
    if len(text.split()) < 150:
        return [text]
    chunks = []
    table_pattern = re.compile(r'(?=Time \| Location)', re.MULTILINE)
    table_sections = table_pattern.split(text)
    non_table = table_sections[0] if table_sections else ''
    tables = table_sections[1:] if len(table_sections) > 1 else []
    if non_table.strip():
        heading_pattern = re.compile(
            r'(?=\n*(?:#{1,3}\s|DAY\s+\d|Day\s+\d))',
            re.MULTILINE | re.IGNORECASE
        )
        sections = heading_pattern.split(non_table)
        sections = [s.strip() for s in sections if s.strip() and len(s.split()) > 10]
        chunks.extend(sections)
    for table in tables:
        if table.strip() and len(table.split()) > 10:
            chunks.append(table.strip())
    if not chunks:
        words = text.split()
        chunk_size = 400
        overlap = 80
        for i in range(0, len(words), chunk_size - overlap):
            chunk = ' '.join(words[i:i + chunk_size])
            if chunk.strip():
                chunks.append(chunk)
    print(f"Chunking: {len(chunks)} chunks for {filename}")
    return chunks

def describe_image_llava(filepath):
    try:
        import requests
        import base64
        with open(filepath, 'rb') as f:
            image_data = base64.b64encode(f.read()).decode('utf-8')
        response = requests.post('http://localhost:11434/api/generate', json={
            'model': 'llava',
            'prompt': 'Describe this image in detail. Include all visible text, objects, people, colors, and any other relevant information.',
            'images': [image_data],
            'stream': False
        }, timeout=60)
        if response.status_code == 200:
            return response.json().get('response', '')
    except Exception as e:
        print(f"LLaVA failed: {e}")
    return None

def embed_and_store(file_id, chunks, user_id, kb_id):
    if not chunks:
        return
    embeddings = embedder.encode(chunks, show_progress_bar=False)
    embeddings = np.array(embeddings, dtype='float32')
    faiss.normalize_L2(embeddings)
    index = faiss.IndexFlatIP(embeddings.shape[1])
    index.add(embeddings)
    folder = f'embeddings/user_{user_id}/kb_{kb_id}'
    os.makedirs(folder, exist_ok=True)
    faiss.write_index(index, f'{folder}/file_{file_id}.index')
    with open(f'{folder}/file_{file_id}.chunks', 'w', encoding='utf-8') as f:
        json.dump(chunks, f, ensure_ascii=False)

def hybrid_search(kb_ids, query, top_k=20, file_ids=None):
    query_vec = embedder.encode([query], show_progress_bar=False)
    query_vec = np.array(query_vec, dtype='float32')
    faiss.normalize_L2(query_vec)
    all_chunks = []
    all_metadata = []

    for kb_id in kb_ids:
        kb = KnowledgeBase.query.get(kb_id)
        if not kb:
            continue
        all_files = get_all_files(kb)
        for f in all_files:
            user_id = kb.user_id
            chunks_path = f'embeddings/user_{user_id}/kb_{kb_id}/file_{f.id}.chunks'
            index_path = f'embeddings/user_{user_id}/kb_{kb_id}/file_{f.id}.index'
            if not os.path.exists(chunks_path):
                if f.content:
                    chunks = semantic_chunk(f.content, f.filename)
                    embed_and_store(f.id, chunks, user_id, kb_id)
                else:
                    continue
            if not os.path.exists(chunks_path):
                continue
            with open(chunks_path, 'r', encoding='utf-8') as cf:
                chunks = json.load(cf)
            for i, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadata.append({
                    'file_id': f.id,
                    'kb_id': kb_id,
                    'filename': f.filename,
                    'kb_name': kb.name,
                    'chunk_idx': i,
                    'index_path': index_path
                })

    if file_ids:
        for file_id in file_ids:
            f = KBFile.query.get(file_id)
            if not f:
                continue
            kb = KnowledgeBase.query.get(f.kb_id)
            if not kb:
                continue
            user_id = kb.user_id
            kb_id = f.kb_id
            chunks_path = f'embeddings/user_{user_id}/kb_{kb_id}/file_{f.id}.chunks'
            index_path = f'embeddings/user_{user_id}/kb_{kb_id}/file_{f.id}.index'
            if not os.path.exists(chunks_path):
                if f.content:
                    chunks = semantic_chunk(f.content, f.filename)
                    embed_and_store(f.id, chunks, user_id, kb_id)
                else:
                    continue
            if not os.path.exists(chunks_path):
                continue
            with open(chunks_path, 'r', encoding='utf-8') as cf:
                chunks = json.load(cf)
            for i, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadata.append({
                    'file_id': f.id,
                    'kb_id': kb_id,
                    'filename': f.filename,
                    'kb_name': kb.name,
                    'chunk_idx': i,
                    'index_path': index_path
                })

    if not all_chunks:
        return []

    tokenized = [c.lower().split() for c in all_chunks]
    bm25 = BM25Okapi(tokenized)
    bm25_scores = bm25.get_scores(query.lower().split())
    faiss_scores = np.zeros(len(all_chunks))
    seen_indexes = {}
    for i, meta in enumerate(all_metadata):
        index_path = meta['index_path']
        if index_path not in seen_indexes:
            if os.path.exists(index_path):
                seen_indexes[index_path] = faiss.read_index(index_path)
        if index_path in seen_indexes:
            idx = seen_indexes[index_path]
            chunk_idx = meta['chunk_idx']
            if chunk_idx < idx.ntotal:
                scores, _ = idx.search(query_vec, idx.ntotal)
                if chunk_idx < len(scores[0]):
                    faiss_scores[i] = max(0, scores[0][chunk_idx])

    bm25_norm = bm25_scores / bm25_scores.max() if bm25_scores.max() > 0 else bm25_scores
    faiss_norm = faiss_scores / faiss_scores.max() if faiss_scores.max() > 0 else faiss_scores
    combined = 0.4 * bm25_norm + 0.6 * faiss_norm
    top_indices = np.argsort(combined)[::-1][:top_k]
    results = []
    for idx in top_indices:
        if combined[idx] > 0:
            results.append({
                'score': float(combined[idx]),
                'content': all_chunks[idx],
                'filename': all_metadata[idx]['filename'],
                'kb_name': all_metadata[idx]['kb_name']
            })
    return results

def extract_text(filepath, filename):
    ext = filename.rsplit('.', 1)[-1].lower()
    try:
        if ext in ['txt', 'md']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext in ['html', 'htm']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                from markdownify import markdownify
                return markdownify(f.read())
        elif ext == 'json':
            import json as jl
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                data = jl.load(f)
                return f"```json\n{jl.dumps(data, indent=2)}\n```"
        elif ext in ['xml', 'csv']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == 'pdf':
            import fitz
            doc = fitz.open(filepath)
            lines = []
            for page in doc:
                try:
                    tabs = page.find_tables()
                    if tabs and tabs.tables:
                        for table in tabs.tables:
                            rows = table.extract()
                            for row in rows:
                                cells = [str(c).strip() if c else '' for c in row]
                                row_text = ' | '.join(c for c in cells if c)
                                if row_text.strip():
                                    lines.append(row_text)
                except:
                    pass
                blocks = page.get_text('dict')['blocks']
                for block in blocks:
                    if block['type'] == 0:
                        for line in block['lines']:
                            text = ' '.join([s['text'] for s in line['spans']]).strip()
                            if not text:
                                continue
                            size = line['spans'][0]['size'] if line['spans'] else 0
                            if size >= 16:
                                lines.append(f"\n# {text}")
                            elif size >= 13:
                                lines.append(f"\n## {text}")
                            else:
                                lines.append(text)
                lines.append('\n')
            return '\n'.join(lines)
        elif ext == 'docx':
            from docx import Document
            doc = Document(filepath)
            lines = []
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name.lower()
                    if 'heading 1' in style:
                        lines.append(f"\n# {para.text}")
                    elif 'heading 2' in style:
                        lines.append(f"\n## {para.text}")
                    elif 'heading 3' in style:
                        lines.append(f"\n### {para.text}")
                    else:
                        lines.append(para.text)
            for table in doc.tables:
                lines.append('\n')
                for row in table.rows:
                    row_cells = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace('\n', ' ')
                        if cell_text:
                            row_cells.append(cell_text)
                    if row_cells:
                        lines.append(' | '.join(row_cells))
                lines.append('\n')
            return '\n'.join(lines)
        elif ext in ['xlsx', 'xls']:
            import openpyxl
            wb = openpyxl.load_workbook(filepath)
            lines = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                lines.append(f"\n## Sheet: {sheet_name}\n")
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                headers = [str(c) if c is not None else '' for c in rows[0]]
                lines.append('| ' + ' | '.join(headers) + ' |')
                lines.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')
                for row in rows[1:]:
                    cells = [str(c) if c is not None else '' for c in row]
                    lines.append('| ' + ' | '.join(cells) + ' |')
            return '\n'.join(lines)
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"\n## Slide {i + 1}\n")
                for shape in slide.shapes:
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            row_text = ' | '.join(c for c in cells if c)
                            if row_text:
                                lines.append(row_text)
                    elif hasattr(shape, 'text') and shape.text.strip():
                        if shape.shape_type == 13:
                            continue
                        lines.append(f"- {shape.text.strip()}")
            return '\n'.join(lines)
        elif ext in ['png', 'jpg', 'jpeg', 'gif', 'webp', 'bmp', 'tiff']:
            description = describe_image_llava(filepath)
            if description:
                return f"[Image: {filename}]\nDescription: {description}"
            return f"[Image file: {filename}]"
        elif ext in ['mp3', 'wav', 'mp4']:
            return f'[Media file: {filename}]'
        else:
            return f'[Unsupported file type: {filename}]'
    except Exception as e:
        return f'[Could not extract text: {str(e)}]'

def get_all_files(kb):
    files = list(kb.files)
    for child in kb.children:
        files.extend(get_all_files(child))
    return files

def get_kb_stats(kb):
    total_size = 0
    last_modified = None
    all_files = get_all_files(kb)
    for f in all_files:
        if os.path.exists(f.filepath):
            total_size += os.path.getsize(f.filepath)
            mtime = os.path.getmtime(f.filepath)
            if last_modified is None or mtime > last_modified:
                last_modified = mtime
    if total_size < 1024:
        size_str = f"{total_size} B"
    elif total_size < 1024 * 1024:
        size_str = f"{total_size / 1024:.1f} KB"
    else:
        size_str = f"{total_size / (1024*1024):.1f} MB"
    date_str = datetime.fromtimestamp(last_modified).strftime('%Y-%m-%d %H:%M') if last_modified else 'No files yet'
    return size_str, date_str, len(all_files)

def flatten_kb_tree(kb, path='', level=0):
    full_path = (path + ' > ' + kb.name) if path else kb.name
    result = [{'id': kb.id, 'name': kb.name, 'path': full_path, 'level': level, 'has_children': len(kb.children) > 0}]
    for child in kb.children:
        result.extend(flatten_kb_tree(child, full_path, level + 1))
    return result

def build_kb_path(kb):
    path_parts = []
    current = kb
    while current:
        path_parts.insert(0, current.name)
        current = current.parent
    return ' > '.join(path_parts)

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SECRET_KEY')
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///knowledge_agent.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db = SQLAlchemy(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

gemini_client = genai.Client(api_key=os.getenv('GEMINI_API_KEY'))
groq_client = Groq(api_key=os.getenv('GROQ_API_KEY'))

@app.template_filter('filesize')
def filesize_filter(filepath):
    try:
        size = os.path.getsize(filepath)
        if size < 1024:
            return f"{size} B"
        elif size < 1024 * 1024:
            return f"{size/1024:.1f} KB"
        else:
            return f"{size/(1024*1024):.1f} MB"
    except:
        return ''

class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password_hash = db.Column(db.String(200), nullable=False)

class Conversation(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(200), default='New Conversation')
    mode = db.Column(db.String(20), default='chat')
    kb_selections = db.Column(db.Text, default='{}')
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    messages = db.relationship('Message', backref='conversation', lazy=True, cascade='all, delete-orphan')

class Message(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    role = db.Column(db.String(10), nullable=False)
    content = db.Column(db.Text, nullable=False)
    conversation_id = db.Column(db.Integer, db.ForeignKey('conversation.id'), nullable=False)

class KnowledgeBase(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    name = db.Column(db.String(200), nullable=False)
    description = db.Column(db.Text, nullable=True, default='')
    instructions = db.Column(db.Text, nullable=True, default='')
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    parent_id = db.Column(db.Integer, db.ForeignKey('knowledge_base.id'), nullable=True)
    children = db.relationship('KnowledgeBase', backref=db.backref('parent', remote_side='KnowledgeBase.id'), lazy=True, cascade='all, delete-orphan')
    files = db.relationship('KBFile', backref='knowledge_base', lazy=True, cascade='all, delete-orphan')

class KBFile(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(300), nullable=False)
    filepath = db.Column(db.String(500), nullable=False)
    content = db.Column(db.Text, nullable=True)
    kb_id = db.Column(db.Integer, db.ForeignKey('knowledge_base.id'), nullable=False)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

@app.route('/', methods=['GET', 'POST'])
def login():
    error = None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password'].encode('utf-8')
        user = User.query.filter_by(username=username).first()
        if user and bcrypt.checkpw(password, user.password_hash):
            login_user(user)
            return redirect(url_for('dashboard'))
        else:
            error = 'Invalid username or password'
    return render_template('login.html', error=error)

@app.route('/register', methods=['GET', 'POST'])
def register():
    error = None
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password'].encode('utf-8')
        confirm_password = request.form['confirm_password'].encode('utf-8')
        if password != confirm_password:
            error = 'Passwords do not match'
        else:
            existing_user = User.query.filter_by(username=username).first()
            if existing_user:
                error = 'Username already taken — please choose another'
            else:
                password_hash = bcrypt.hashpw(password, bcrypt.gensalt())
                new_user = User(username=username, password_hash=password_hash)
                db.session.add(new_user)
                db.session.commit()
                return redirect(url_for('login'))
    return render_template('register.html', error=error)

@app.route('/dashboard')
@app.route('/dashboard/<int:conversation_id>')
@login_required
def dashboard(conversation_id=None):
    conversations = Conversation.query.filter_by(user_id=current_user.id).all()
    active_conversation = None
    messages = []
    if conversation_id:
        active_conversation = Conversation.query.get(conversation_id)
    elif conversations:
        active_conversation = conversations[-1]
    if active_conversation:
        messages = Message.query.filter_by(conversation_id=active_conversation.id).all()
    kbs = KnowledgeBase.query.filter_by(user_id=current_user.id).all()
    kb_stats = {kb.id: get_kb_stats(kb) for kb in kbs}
    root_kbs = [kb for kb in kbs if not kb.parent_id]
    kb_tree_flat = []
    for kb in root_kbs:
        kb_tree_flat.extend(flatten_kb_tree(kb))
    active_kb_selections = []
    if active_conversation and active_conversation.kb_selections:
        try:
            selections = json.loads(active_conversation.kb_selections)
            if isinstance(selections, list):
                kb_ids_sel = selections
                file_ids_sel = []
            else:
                kb_ids_sel = selections.get('kb_ids', [])
                file_ids_sel = selections.get('file_ids', [])
            for kb_id in kb_ids_sel:
                kb = KnowledgeBase.query.get(kb_id)
                if kb:
                    active_kb_selections.append({'id': kb.id, 'path': build_kb_path(kb)})
            for file_id in file_ids_sel:
                f = KBFile.query.get(file_id)
                if f:
                    kb = KnowledgeBase.query.get(f.kb_id)
                    if kb:
                        active_kb_selections.append({'id': f'file_{f.id}', 'path': build_kb_path(kb) + ' > ' + f.filename})
        except:
            pass
    return render_template('dashboard.html',
        conversations=conversations,
        active_conversation=active_conversation,
        messages=messages,
        kbs=kbs,
        root_kbs=root_kbs,
        kb_tree_flat=kb_tree_flat,
        kb_stats=kb_stats,
        active_kb_selections=active_kb_selections)

@app.route('/conversation/new', methods=['GET', 'POST'])
@login_required
def new_conversation():
    if request.method == 'POST':
        data = request.get_json()
        title = data.get('title', 'New Conversation')
        mode = data.get('mode', 'chat')
        kb_ids = data.get('kb_ids', [])
        file_ids = data.get('file_ids', [])
        selections = {'kb_ids': kb_ids, 'file_ids': file_ids}
        conv = Conversation(title=title, mode=mode, kb_selections=json.dumps(selections), user_id=current_user.id)
        db.session.add(conv)
        db.session.commit()
        return jsonify({'conversation_id': conv.id})
    conv = Conversation(title='New Conversation', user_id=current_user.id)
    db.session.add(conv)
    db.session.commit()
    return redirect(url_for('dashboard', conversation_id=conv.id))

@app.route('/conversation/<int:conversation_id>/update', methods=['POST'])
@login_required
def update_conversation(conversation_id):
    conv = Conversation.query.get(conversation_id)
    if not conv or conv.user_id != current_user.id:
        return jsonify({'error': 'Invalid'}), 400
    data = request.get_json()
    if 'title' in data:
        conv.title = data['title']
    if 'mode' in data:
        conv.mode = data['mode']
    if 'kb_ids' in data:
        try:
            existing = json.loads(conv.kb_selections or '{}')
            if isinstance(existing, list):
                existing = {'kb_ids': existing, 'file_ids': []}
        except:
            existing = {'kb_ids': [], 'file_ids': []}
        existing['kb_ids'] = data['kb_ids']
        conv.kb_selections = json.dumps(existing)
    db.session.commit()
    return jsonify({'success': True})

@app.route('/conversation/<int:conversation_id>/delete')
@login_required
def delete_conversation(conversation_id):
    conv = Conversation.query.get(conversation_id)
    if conv and conv.user_id == current_user.id:
        db.session.delete(conv)
        db.session.commit()
    return redirect(url_for('dashboard'))

def call_gemini(message):
    response = gemini_client.models.generate_content(model='gemini-2.5-flash', contents=message)
    return response.text

def call_groq(message):
    response = groq_client.chat.completions.create(
        model='llama-3.3-70b-versatile',
        messages=[{'role': 'user', 'content': message}]
    )
    return response.choices[0].message.content

def call_ollama(message):
    import requests
    response = requests.post('http://localhost:11434/api/chat', json={
        'model': 'llama3.2',
        'messages': [{'role': 'user', 'content': message}],
        'stream': False
    })
    return response.json()['message']['content']

def call_gemma(message):
    import requests
    response = requests.post('http://localhost:11434/api/chat', json={
        'model': 'gemma3:4b',
        'messages': [{'role': 'user', 'content': message}],
        'stream': False
    })
    return response.json()['message']['content']

@app.route('/chat', methods=['POST'])
@login_required
def chat():
    data = request.get_json()
    user_message = data.get('message', '')
    conversation_id = data.get('conversation_id')
    model = data.get('model', 'gemini')

    conv = Conversation.query.get(conversation_id)
    if not conv or conv.user_id != current_user.id:
        return jsonify({'error': 'Invalid conversation'}), 400

    if conv.title == 'New Conversation':
        conv.title = user_message[:50]
        db.session.commit()

    user_msg = Message(role='user', content=user_message, conversation_id=conv.id)
    db.session.add(user_msg)
    db.session.commit()

    prompt = user_message
    sources_used = []

    if conv.mode == 'knowledge':
        try:
            selections = json.loads(conv.kb_selections or '{}')
            if isinstance(selections, list):
                kb_ids = selections
                file_ids = []
            else:
                kb_ids = selections.get('kb_ids', [])
                file_ids = selections.get('file_ids', [])
        except:
            kb_ids = []
            file_ids = []

        if kb_ids or file_ids:
            chunks = hybrid_search(kb_ids, user_message, top_k=20, file_ids=file_ids)
            if chunks:
                context_parts = []
                total_words = 0
                token_budget = 20000
                for chunk in chunks:
                    chunk_words = len(chunk['content'].split())
                    if total_words + chunk_words > token_budget:
                        break
                    context_parts.append(f"[{chunk['kb_name']} / {chunk['filename']}]\n{chunk['content']}")
                    sources_used.append(f"{chunk['kb_name']} / {chunk['filename']}")
                    total_words += chunk_words

                instructions_parts = []
                for kb_id in kb_ids:
                    kb = KnowledgeBase.query.get(kb_id)
                    if kb and kb.instructions:
                        instructions_parts.append(f"Instructions for {kb.name}: {kb.instructions}")

                context_block = '\n\n---\n\n'.join(context_parts)
                instructions_section = f"INSTRUCTIONS:\n{chr(10).join(instructions_parts)}\n\n" if instructions_parts else ""

                prompt = f"""You are a smart friendly personal assistant with access to the user's personal documents.

{instructions_section}RELEVANT CONTENT FROM USER'S DOCUMENTS:
{context_block}

Guidelines:
- Be natural, warm and conversational
- Use the document content above as your primary source
- Give complete detailed answers using ALL relevant information
- When documents contain tables with times/locations/activities, list them all clearly
- Combine document knowledge with general knowledge when helpful
- For greetings or off-topic questions just respond normally

User: {user_message}
Assistant:"""

    try:
        if model == 'groq':
            answer = call_groq(prompt)
        elif model == 'ollama':
            answer = call_ollama(prompt)
        elif model == 'gemma':
            answer = call_gemma(prompt)
        else:
            answer = call_gemini(prompt)
    except Exception as e:
        try:
            answer = call_groq(prompt)
            answer = '⚡ (Gemini unavailable, using Groq as backup)\n\n' + answer
        except:
            return jsonify({'error': 'All models unavailable. Please try again.'}), 503

    if sources_used:
        unique_sources = list(dict.fromkeys(sources_used))
        answer += '\n\n---\n📚 **Sources:** ' + ', '.join(unique_sources)

    agent_msg = Message(role='agent', content=answer, conversation_id=conv.id)
    db.session.add(agent_msg)
    db.session.commit()

    actual_model = model
    if model == 'gemini' and '⚡' in answer:
        actual_model = 'groq'
    return jsonify({'response': answer, 'conversation_id': conv.id, 'model_used': actual_model})

@app.route('/logout')
@login_required
def logout():
    logout_user()
    return redirect(url_for('login'))

@app.route('/kb')
@login_required
def kb_manager():
    return redirect(url_for('dashboard') + '?mode=kb')

@app.route('/kb/new', methods=['POST'])
@login_required
def kb_new():
    name = request.form.get('name', '').strip()
    parent_id = request.form.get('parent_id', None)
    if parent_id:
        parent_id = int(parent_id)
    if name:
        kb = KnowledgeBase(name=name, user_id=current_user.id, parent_id=parent_id)
        db.session.add(kb)
        db.session.commit()
        os.makedirs(f'uploads/user_{current_user.id}/kb_{kb.id}', exist_ok=True)
    if parent_id:
        return redirect(url_for('kb_detail', kb_id=parent_id))
    return redirect(url_for('dashboard') + '?mode=kb&view=manage')

@app.route('/kb/import', methods=['POST'])
@login_required
def kb_import_global():
    from werkzeug.utils import secure_filename
    file = request.files.get('zip_file')
    if not file or not file.filename.endswith('.zip'):
        return redirect(url_for('dashboard') + '?mode=kb')
    zip_name = file.filename.replace('.zip', '')
    new_kb = KnowledgeBase(name=zip_name, user_id=current_user.id, parent_id=None)
    db.session.add(new_kb)
    db.session.commit()
    folder = f'uploads/user_{current_user.id}/kb_{new_kb.id}'
    os.makedirs(folder, exist_ok=True)
    zip_buffer = io.BytesIO(file.read())
    with zipfile.ZipFile(zip_buffer, 'r') as zip_file:
        for zip_entry in zip_file.namelist():
            if zip_entry.endswith('/'):
                continue
            filename = secure_filename(os.path.basename(zip_entry))
            if not filename:
                continue
            filepath = os.path.join(folder, filename)
            with open(filepath, 'wb') as f:
                f.write(zip_file.read(zip_entry))
            content = extract_text(filepath, filename)
            kb_file = KBFile(filename=filename, filepath=filepath, content=content, kb_id=new_kb.id)
            db.session.add(kb_file)
            db.session.flush()
            chunks = semantic_chunk(content, filename)
            embed_and_store(kb_file.id, chunks, current_user.id, new_kb.id)
    db.session.commit()
    return redirect(url_for('kb_detail', kb_id=new_kb.id))

@app.route('/kb/<int:kb_id>')
@login_required
def kb_detail(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return redirect(url_for('dashboard'))
    breadcrumb = []
    parent = kb.parent
    while parent:
        breadcrumb.insert(0, parent)
        parent = parent.parent
    return render_template('kb_detail.html', kb=kb, breadcrumb=breadcrumb)

@app.route('/kb/<int:kb_id>/explore')
@login_required
def kb_explore(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return redirect(url_for('dashboard'))
    breadcrumb = []
    parent = kb.parent
    while parent:
        breadcrumb.insert(0, parent)
        parent = parent.parent
    return render_template('kb_explore.html', kb=kb, breadcrumb=breadcrumb)

@app.route('/kb/<int:kb_id>/delete')
@login_required
def kb_delete(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    parent_id = kb.parent_id if kb else None
    if kb and kb.user_id == current_user.id:
        db.session.delete(kb)
        db.session.commit()
    if parent_id:
        return redirect(url_for('kb_detail', kb_id=parent_id))
    return redirect(url_for('dashboard') + '?mode=kb&view=manage')

@app.route('/kb/<int:kb_id>/rename', methods=['POST'])
@login_required
def kb_rename(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if kb and kb.user_id == current_user.id:
        new_name = request.form.get('name', '').strip()
        if new_name:
            kb.name = new_name
            db.session.commit()
    if kb and kb.parent_id:
        return redirect(url_for('kb_detail', kb_id=kb.parent_id))
    return redirect(url_for('kb_detail', kb_id=kb_id))

@app.route('/kb/<int:kb_id>/upload', methods=['POST'])
@login_required
def kb_upload(kb_id):
    from werkzeug.utils import secure_filename
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return redirect(url_for('dashboard') + '?mode=kb')
    file = request.files.get('file')
    if file and file.filename:
        filename = secure_filename(file.filename)
        folder = f'uploads/user_{current_user.id}/kb_{kb_id}'
        os.makedirs(folder, exist_ok=True)
        filepath = os.path.join(folder, filename)
        file.save(filepath)
        content = extract_text(filepath, filename)
        kb_file = KBFile(filename=filename, filepath=filepath, content=content, kb_id=kb_id)
        db.session.add(kb_file)
        db.session.flush()
        chunks = semantic_chunk(content, filename)
        embed_and_store(kb_file.id, chunks, current_user.id, kb_id)
        db.session.commit()
    return redirect(url_for('kb_detail', kb_id=kb_id))

@app.route('/kb/<int:kb_id>/file/<int:file_id>/delete')
@login_required
def kb_file_delete(kb_id, file_id):
    f = KBFile.query.get(file_id)
    if f and f.kb_id == kb_id:
        if os.path.exists(f.filepath):
            os.remove(f.filepath)
        kb = KnowledgeBase.query.get(kb_id)
        if kb:
            for ext in ['index', 'chunks']:
                p = f'embeddings/user_{kb.user_id}/kb_{kb_id}/file_{file_id}.{ext}'
                if os.path.exists(p):
                    os.remove(p)
        db.session.delete(f)
        db.session.commit()
    return redirect(url_for('kb_detail', kb_id=kb_id))

@app.route('/kb/<int:kb_id>/file/<int:file_id>/rename', methods=['POST'])
@login_required
def kb_file_rename(kb_id, file_id):
    f = KBFile.query.get(file_id)
    if not f or f.kb_id != kb_id:
        return jsonify({'error': 'Invalid file'}), 400
    new_name = request.form.get('new_name', '').strip()
    if not new_name:
        return jsonify({'error': 'Invalid name'}), 400
    old_path = f.filepath
    folder = os.path.dirname(old_path)
    new_path = os.path.join(folder, new_name)
    if os.path.exists(old_path):
        os.rename(old_path, new_path)
    f.filename = new_name
    f.filepath = new_path
    db.session.commit()
    return jsonify({'success': True})

@app.route('/kb/<int:kb_id>/description/save', methods=['POST'])
@login_required
def kb_description_save(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return jsonify({'success': False}), 400
    kb.description = request.form.get('description', '').strip()
    db.session.commit()
    return redirect(url_for('kb_detail', kb_id=kb_id))

@app.route('/kb/<int:kb_id>/instructions/save', methods=['POST'])
@login_required
def kb_instructions_save(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return jsonify({'success': False}), 400
    kb.instructions = request.form.get('instructions', '').strip()
    db.session.commit()
    return redirect(url_for('kb_detail', kb_id=kb_id))

@app.route('/kb/<int:kb_id>/export/selected', methods=['POST'])
@login_required
def kb_export_selected(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return jsonify({'error': 'Invalid KB'}), 400
    file_ids = request.json.get('file_ids', [])
    if not file_ids:
        return jsonify({'error': 'No files selected'}), 400
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for file_id in file_ids:
            kb_file = KBFile.query.get(int(file_id))
            if kb_file and kb_file.kb_id == kb_id and os.path.exists(kb_file.filepath):
                zip_file.write(kb_file.filepath, arcname=kb_file.filename)
    zip_buffer.seek(0)
    now = datetime.now().strftime('%Y-%m-%d')
    return send_file(zip_buffer, mimetype='application/zip', as_attachment=True, download_name=f'{kb.name}_{now}.zip')

@app.route('/kb/<int:kb_id>/export/all')
@login_required
def kb_export_all(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return redirect(url_for('dashboard') + '?mode=kb')
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for kb_file in kb.files:
            if os.path.exists(kb_file.filepath):
                zip_file.write(kb_file.filepath, arcname=kb_file.filename)
    zip_buffer.seek(0)
    now = datetime.now().strftime('%Y-%m-%d')
    return send_file(zip_buffer, mimetype='application/zip', as_attachment=True, download_name=f'{kb.name}_{now}.zip')

@app.route('/kb/<int:kb_id>/files/delete-selected', methods=['POST'])
@login_required
def kb_files_delete_selected(kb_id):
    kb = KnowledgeBase.query.get(kb_id)
    if not kb or kb.user_id != current_user.id:
        return redirect(url_for('dashboard') + '?mode=kb')
    file_ids = request.json.get('file_ids', [])
    for file_id in file_ids:
        f = KBFile.query.get(int(file_id))
        if f and f.kb_id == kb_id:
            if os.path.exists(f.filepath):
                os.remove(f.filepath)
            for ext in ['index', 'chunks']:
                p = f'embeddings/user_{kb.user_id}/kb_{kb_id}/file_{f.id}.{ext}'
                if os.path.exists(p):
                    os.remove(p)
            db.session.delete(f)
    db.session.commit()
    return jsonify({'success': True})

if __name__ == '__main__':
    with app.app_context():
        db.create_all()
    app.run(debug=True)